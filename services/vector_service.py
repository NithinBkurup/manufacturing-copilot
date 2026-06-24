"""
Vector Service – Manufacturing Copilot
ChromaDB-backed RAG for SOPs, work instructions, recipes, and manuals.
Supported file types: PDF, DOCX, XLSX, PPTX, TXT
"""

import logging
import os
from pathlib import Path
from typing import List, Optional

from config.settings import settings

logger = logging.getLogger("copilot.vector")

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100


class VectorService:

    def __init__(self):
        self._client = None
        self._collection = None
        self._embed_fn = None

    async def initialize(self) -> None:
        try:
            import chromadb
            import os
            from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction

            os.makedirs(settings.CHROMA_PERSIST_DIR, exist_ok=True)
            os.environ["TRANSFORMERS_OFFLINE"] = "1"
            os.environ["HF_DATASETS_OFFLINE"] = "1"
            cache_dir = "D:\\Dev\\sentence_transformers_cache"
            self._client = chromadb.PersistentClient(path=settings.CHROMA_PERSIST_DIR)
            self._embed_fn = SentenceTransformerEmbeddingFunction(
                model_name="all-MiniLM-L6-v2",
                cache_folder=cache_dir,
            )
            self._collection = self._client.get_or_create_collection(
                name=settings.CHROMA_COLLECTION_DOCS,
                embedding_function=self._embed_fn,
            )
            logger.info(
                "ChromaDB initialised — collection '%s' has %d documents",
                settings.CHROMA_COLLECTION_DOCS,
                self._collection.count(),
            )
        except Exception as exc:
            logger.warning("ChromaDB init failed (RAG disabled): %s", exc)

    def get_context_for_ai(self, query: str, top_k: Optional[int] = None) -> str:
        if self._collection is None:
            return ""
        try:
            k = top_k or settings.CHROMA_TOP_K
            results = self._collection.query(query_texts=[query], n_results=k)
            if not results["documents"] or not results["documents"][0]:
                return ""

            parts = ["[KNOWLEDGE BASE — relevant document excerpts]"]
            for i, (doc, meta) in enumerate(
                zip(results["documents"][0], results["metadatas"][0]), start=1
            ):
                source = meta.get("source", "Unknown Document")
                parts.append(f"\n[{i}] Source: {source}\n{doc}")
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("RAG query error: %s", exc)
            return ""

    async def index_document(self, file_path: str, original_name: Optional[str] = None) -> int:
        """Index a document file or directory recursively into ChromaDB. Returns number of chunks added."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")

        if path.is_dir():
            total_chunks = 0
            indexed_any = False
            supported_suffixes = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt", ".txt"}
            files_to_index = []
            for root, _, filenames in os.walk(path):
                for f in filenames:
                    fp = Path(root) / f
                    if fp.suffix.lower() in supported_suffixes:
                        files_to_index.append(fp)

            if not files_to_index:
                raise ValueError("No supported files found in the directory.")

            already_indexed_count = 0
            for fp in files_to_index:
                # Check if already indexed
                if self._collection is not None:
                    existing = self._collection.get(where={"source": fp.name})
                    if existing and existing.get("ids"):
                        already_indexed_count += 1
                        continue  # Skip already indexed file in directory walk
                
                try:
                    chunks = await self._index_single_file(fp, fp.name)
                    if chunks > 0:
                        indexed_any = True
                        total_chunks += chunks
                except Exception as e:
                    logger.error("Failed to index %s during walk: %s", fp, e)
            
            if not indexed_any and already_indexed_count > 0:
                raise ValueError("All files in this directory are already indexed.")
            return total_chunks

        else:
            source_name = original_name or path.name
            if self._collection is not None:
                existing = self._collection.get(where={"source": source_name})
                if existing and existing.get("ids"):
                    raise ValueError(f"File '{source_name}' is already indexed in the Knowledge Base.")
            
            return await self._index_single_file(path, source_name)

    async def _index_single_file(self, path: Path, source_name: str) -> int:
        text = self._extract_text(path)
        if not text:
            logger.warning("No text extracted from %s", path)
            return 0

        chunk_size = getattr(settings, "KB_CHUNK_SIZE", 800)
        chunk_overlap = getattr(settings, "KB_CHUNK_OVERLAP", 100)
        chunks = self._chunk_text(text, chunk_size, chunk_overlap)
        if not chunks:
            return 0

        ids = [f"{path.stem}_{i}_{abs(hash(source_name)) & 0xffffffff}" for i in range(len(chunks))]
        metadatas = [{"source": source_name, "file_path": str(path)} for _ in chunks]

        self._collection.upsert(ids=ids, documents=chunks, metadatas=metadatas)
        logger.info("Indexed %d chunks from %s", len(chunks), source_name)
        return len(chunks)

    def _extract_text(self, path: Path) -> str:
        suffix = path.suffix.lower()
        try:
            if suffix == ".txt":
                return path.read_text(encoding="utf-8", errors="ignore")
            elif suffix == ".pdf":
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            elif suffix == ".docx":
                from docx import Document
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs)
            elif suffix in (".xlsx", ".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        rows.append("\t".join(str(c) if c is not None else "" for c in row))
                return "\n".join(rows)
            elif suffix in (".pptx", ".ppt"):
                from pptx import Presentation
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
        except Exception as exc:
            logger.error("Text extraction failed for %s: %s", path.name, exc)
        return ""

    def _chunk_text(self, text: str, chunk_size: int = 800, chunk_overlap: int = 100) -> List[str]:
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - chunk_overlap
        return [c.strip() for c in chunks if c.strip()]

    @property
    def indexed_files(self) -> List[dict]:
        if self._collection is None:
            return []
        try:
            results = self._collection.get(include=["metadatas"])
            if not results or not results.get("metadatas"):
                return []
            seen = set()
            files = []
            for meta in results["metadatas"]:
                if not meta:
                    continue
                source = meta.get("source")
                path = meta.get("file_path")
                if source and source not in seen:
                    seen.add(source)
                    files.append({"source": source, "file_path": path or ""})
            return files
        except Exception as exc:
            logger.warning("Failed to get indexed files: %s", exc)
            return []

    @property
    def document_count(self) -> int:
        if self._collection is None:
            return 0
        return self._collection.count()